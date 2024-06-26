from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import copy

with open("cantico.txt", "r", encoding="utf8") as file:
  lines = file.readlines()


def getCanticoByLineInText(initalLineCantico, finalLineCantico):
  cantico = []
  for line in lines[initalLineCantico:finalLineCantico]:
      cantico.append(line)
  return cantico

def getPositionLine(startCantico, endCantito):
  #Pegar linha com numero do inio desejado
  initalLineCantico = 0
  finalLineCantico = 0
  for line in lines:
      if startCantico in line:
          print(lines.index(line))
          initalLineCantico = lines.index(line)
      if endCantito in line:
          print(lines.index(line))
          finalLineCantico = lines.index(line)
          break
  return initalLineCantico, finalLineCantico

def getCanticoByNumber(canticoNumber):
  quantityCanticos = canticoNumber
  begginRange = canticoNumber
  for i in range(begginRange,quantityCanticos+1):
    return  getPositionLine(str(i), str(i+1))

def getCantico(numberCantico):
  initalLineCantico, finalLineCantico = getCanticoByNumber(numberCantico)
  return getCanticoByLineInText(initalLineCantico, finalLineCantico)

def printEstrofe(line, ToLine):
   text= cantico[line:ToLine]
   text1 = ""
   for line in text:
     text1 += line
   return text1

def replace_space_after_number(text):
  index = text.find(' ')  # Encontra o índice do primeiro espaço
  if index > 0 and text[:index].isdigit():  # Verifica se há um número antes do espaço
     modified_text = text[:index] + '. ' + text[index + 1:]  # Substitui o espaço por um ponto e espaço
  return modified_text.replace("\n","")
  
def makeSlide(objectCantico):

  prs = Presentation('modelo.pptx')

  for i in range(len(objectCantico)):
    if i == 0:
      #Slide Title
      slide = prs.slides[i]
      shapes =  slide.shapes
      shape = shapes[1]

      text_frame = shape.text_frame

      # Manter a cor e tamanho da fonte do texto original
      text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 102, 0)
      text_frame.paragraphs[0].runs[0].text = objectCantico["texto"+str(i)]

    if (i > 0 and i <= 1):
      #Slide
      slide = prs.slides[i]
      shapes = slide.shapes
      shape = shapes[3]

      text_frame = shape.text_frame
      text_frame.paragraphs[1].runs[0].text = objectCantico.get("texto"+str(i))

    if i > 1:
      # create new slide,
      newSlide = prs.slides.add_slide(prs.slide_layouts[0])
      for shp in slide.shapes:
              el = shp.element
              newel = copy.deepcopy(el)
              newSlide.shapes._spTree.insert_element_before(newel, 'p:extLst')
      #slide
      shapesnw = newSlide.shapes
      shape4 = shapesnw[3]

      text_frame4 = shape4.text_frame
      text_frame4.paragraphs[1].runs[0].text = objectCantico.get("texto"+str(i))

  prs.save("./hinos/"+ objectCantico["texto"+str(0)]+'.pptx')

 

def makeObjectCantico(posicaoIndexCaracterNewLine):
  objectCantico = {}

  for  i in range(len(posicaoIndexCaracterNewLine)):
    if i == len(posicaoIndexCaracterNewLine)-1:
      break
    objectCantico["texto"+ str(i)] = printEstrofe(posicaoIndexCaracterNewLine[i], posicaoIndexCaracterNewLine[i+1])


  objectCantico["texto0"] = replace_space_after_number(objectCantico["texto0"]).upper()
  return objectCantico

def getPositionCaracterNewLine(cantico):
  firstPositionTitle = 0
  posicaoIndexCaracterNewLine = [firstPositionTitle,]
  pnultimoValor=1
  for i in range(len(cantico)):
    if len(cantico[i]) == 1:
      if (pnultimoValor == i-1):
          # print("É:"+  str(i) + " Parou! :)")
          break
      posicaoIndexCaracterNewLine.append(i)
      pnultimoValor = i
  return(posicaoIndexCaracterNewLine)


#---------Main-----------#
# [110, 237, 281, 289, 325, 354, 377, 378, 400]
# erros = [] 
# for i in range(1,401):
#    try:
#     cantico = getCantico(i)

#     positionCaracterNewLine = getPositionCaracterNewLine(cantico)

#     dictOfCanticos = makeObjectCantico(positionCaracterNewLine)

#     makeSlide(dictOfCanticos)
#    except:
#       print(erros.append(i))
  
# print(erros)

cantico = getCantico(11)

positionCaracterNewLine = getPositionCaracterNewLine(cantico)

dictOfCanticos = makeObjectCantico(positionCaracterNewLine)

makeSlide(dictOfCanticos)